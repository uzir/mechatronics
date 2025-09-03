import streamlit as st
import google.generativeai as genai
import PyPDF2
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import io

# --- פונקציות עזר ---

@st.cache_data
def load_knowledge_base(file_path):
    """קוראת תוכן מקובץ PDF ומחזירה אותו כטקסט."""
    try:
        with open(file_path, "rb") as pdf_file:
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text()
            return text
    except FileNotFoundError:
        return None
    except Exception as e:
        st.error(f"שגיאה בקריאת קובץ ה-PDF: {e}")
        return None

def create_presentation_from_text(text_content):
    """יוצרת מצגת PowerPoint על בסיס טקסט בפורמט Markdown."""
    prs = Presentation()
    slides_text = text_content.strip().split("\n\n")

    TITLE_AND_CONTENT_LAYOUT = 1
    TITLE_ONLY_LAYOUT = 5

    for slide_text in slides_text:
        lines = slide_text.strip().split('\n')
        if not lines or not lines[0]: continue

        title = lines[0].replace("#", "").strip()
        content_points = [line.replace("-", "").strip() for line in lines[1:] if line.strip().startswith("-")]

        if content_points:
            slide_layout = prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT]
            slide = prs.slides.add_slide(slide_layout)
            content_shape = slide.shapes.placeholders[1]
            tf = content_shape.text_frame
            tf.clear()
            for point in content_points:
                p = tf.add_paragraph()
                p.text = point
                p.alignment = PP_ALIGN.RIGHT
                p.font.size = Pt(20)
                p.level = 0
        else:
            slide_layout = prs.slide_layouts[TITLE_ONLY_LAYOUT]
            slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.alignment = PP_ALIGN.RIGHT
        title_paragraph.font.size = Pt(36)
            
    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()

def get_response_text(response):
    """
    פונקציה חסינה לחילוץ טקסט מאובייקט התגובה של המודל.
    מנסה את הגישה הסטנדרטית (.text) ואם נכשלת, מנסה גישת גיבוי.
    """
    try:
        # הדרך הסטנדרטית והמועדפת
        return response.text
    except Exception:
        # דרך גיבוי למקרה של מבנה אובייקט שונה או גרסה ישנה
        try:
            return response.parts[0].text
        except Exception as e:
            st.error(f"שגיאה קריטית: לא ניתן לחלץ טקסט מאובייקט התגובה. {e}")
            st.write(response) # הצגת האובייקט המלא לצורכי דיבאג
            return "שגיאה בהצגת התגובה."

# --- הגדרות והוראות לבוט ---

knowledge_base_text = load_knowledge_base("819387ALL.pdf")
BASE_SYSTEM_INSTRUCTION = """
אתה מורה מומחה במגמות מכטרוניקה (כיתות י–י"ב) עם שלושה מצבים...
(ההנחיות שלך נשארות כפי שהן)
"""
if knowledge_base_text:
    SYSTEM_INSTRUCTION = f"{BASE_SYSTEM_INSTRUCTION}\n\n---מאגר ידע קבוע---\n{knowledge_base_text}\n---"
else:
    SYSTEM_INSTRUCTION = BASE_SYSTEM_INSTRUCTION

PAGE_TITLE = "🤖 המורה למכטרוניקה"
st.set_page_config(page_title=PAGE_TITLE, page_icon="🤖", layout="wide")
st.markdown("""<style> body { direction: rtl; } .stTextInput > div > div > input { direction: rtl; } </style>""", unsafe_allow_html=True)
st.title(PAGE_TITLE)

try:
    genai.configure(api_key=st.secrets["GOOGLE_API_KEY"])
    chat_model = genai.GenerativeModel(model_name="gemini-1.5-pro-latest", system_instruction=SYSTEM_INSTRUCTION, tools=['google_search_retrieval'])
    basic_model = genai.GenerativeModel(model_name="gemini-1.5-pro-latest")
except Exception as e:
    st.error("שגיאה בהגדרת ה-API Key. אנא ודא שהוספת אותו כראוי ב'סודות' האפליקציה.", icon="🚨")
    st.stop()

tabs = st.tabs(["💬 צ'אט עם הבוט", "🖼️ ניתוח תמונות", "🧠 מחולל מבחנים", "📊 מחולל מצגות"])

# --- טאב 1: צ'אט רגיל ---
with tabs[0]:
    st.header("שיחה עם המורה למכטרוניקה")
    # (קוד הצ'אט נשאר ללא שינוי)
    if "chat" not in st.session_state:
        st.session_state.chat = chat_model.start_chat(history=[])
        st.session_state.messages = [{"role": "assistant", "content": "שלום, אני המורה הדיגיטלי למכטרוניקה. איך אוכל לעזור?"}]
    for message in st.session_state.get("messages", []):
        with st.chat_message(message["role"]):
            st.markdown(f'<div style="direction: rtl;">{message["content"]}</div>', unsafe_allow_html=True)
    if prompt := st.chat_input("כתבו כאן את שאלתכם..."):
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(f'<div style="direction: rtl;">{prompt}</div>', unsafe_allow_html=True)
        with st.chat_message("assistant"):
            with st.spinner("חושב..."):
                response_stream = st.session_state.chat.send_message(prompt, stream=True)
                full_response = st.write_stream(response_stream)
        st.session_state.messages.append({"role": "assistant", "content": full_response})

# --- טאב 2: ניתוח תמונות ---
with tabs[1]:
    st.header("ניתוח שרטוטים ותמונות")
    uploaded_image = st.file_uploader("בחר קובץ תמונה", type=["png", "jpg", "jpeg"], key="analyzer")
    image_prompt = st.text_input("מה תרצה לשאול על התמונה?", key="image_q")
    if st.button("נתח את התמונה", key="analyze_btn"):
        if uploaded_image and image_prompt:
            with st.spinner("מעבד את התמונה..."):
                image_obj = Image.open(uploaded_image)
                response = basic_model.generate_content([image_prompt, image_obj])
                st.subheader("תוצאות הניתוח:")
                # שימוש בפונקציה החסינה לחילוץ הטקסט
                response_text = get_response_text(response)
                st.markdown(f'<div style="direction: rtl;">{response_text}</div>', unsafe_allow_html=True)
        else:
            st.warning("אנא העלה תמונה וכתוב שאלה.")

# --- טאב 3: מחולל מבחנים ---
with tabs[2]:
    st.header("מחולל מבחנים וחידונים")
    with st.form("quiz_form"):
        quiz_topic = st.text_input("נושא המבחן")
        num_questions = st.number_input("מספר שאלות", min_value=1, max_value=20, value=5)
        question_type = st.selectbox("סוג השאלות", ["רב-ברירתיות (אמריקאיות)", "פתוחות", "נכון / לא נכון"])
        submitted = st.form_submit_button("🚀 צור את המבחן")
    if submitted and quiz_topic:
        with st.spinner("מכין מבחן..."):
            quiz_prompt = f"צור מבחן בנושא {quiz_topic}, עם {num_questions} שאלות מסוג {question_type}."
            response = basic_model.generate_content(quiz_prompt)
            st.subheader(f"מבחן בנושא: {quiz_topic}")
            # שימוש בפונקציה החסינה לחילוץ הטקסט
            response_text = get_response_text(response)
            st.markdown(f'<div style="direction: rtl;">{response_text}</div>', unsafe_allow_html=True)

# --- טאב 4: מחולל מצגות ---
with tabs[3]:
    st.header("מחולל מצגות PowerPoint")
    with st.form("ppt_form"):
        ppt_topic = st.text_input("נושא המצגת")
        slide_count = st.number_input("מספר שקופיות", min_value=3, max_value=20, value=7)
        target_audience = st.text_input("קהל יעד")
        submitted = st.form_submit_button("📊 צור מצגת")
    if submitted and ppt_topic:
        with st.spinner("כותב את תוכן המצגת..."):
            ppt_prompt = f"צור תוכן למצגת של {slide_count} שקופיות בנושא '{ppt_topic}' לקהל יעד של '{target_audience}'. החזר בפורמט Markdown..."
            response = basic_model.generate_content(ppt_prompt)
            # שימוש בפונקציה החסינה לחילוץ הטקסט
            presentation_text = get_response_text(response)
        with st.spinner("בונה את קובץ ה-PowerPoint..."):
            ppt_file_data = create_presentation_from_text(presentation_text)
        st.success("המצגת שלך מוכנה להורדה!")
        st.download_button(label="📥 הורד את המצגת (.pptx)", data=ppt_file_data, file_name=f"{ppt_topic}.pptx")
        with st.expander("הצג את התוכן הטקסטואלי של המצגת"):
            st.markdown(f'<div style="direction: rtl; text-align: right;">{presentation_text}</div>', unsafe_allow_html=True)
    elif submitted:
        st.warning("אנא מלא את נושא המצגת.")
